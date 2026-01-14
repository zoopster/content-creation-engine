"""
PPTX Generation Skill - Creates PowerPoint presentations from draft content.

This skill transforms markdown content into formatted PPTX files with:
- Professional slide layouts
- Brand template support
- Title, content, and section slides
- Bullet points and formatted text
"""

import sys
import os
sys.path.append(os.path.join(os.path.dirname(__file__), '../..'))

from typing import Dict, Any, Optional, List
from agents.base.agent import Skill
from agents.base.models import DraftContent
from datetime import datetime
import re


class PptxGenerationSkill(Skill):
    """
    Generates PowerPoint (PPTX) presentations from draft content.

    Features:
    - Markdown to PPTX conversion
    - Brand template application
    - Multiple slide layouts (title, content, section)
    - Bullet point support
    - Brand colors and typography
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        super().__init__("pptx-generation", config)
        self.output_dir = config.get("output_dir", "./output") if config else "./output"

        # Ensure directories exist
        os.makedirs(self.output_dir, exist_ok=True)

    def execute(self, input_data: DraftContent, **kwargs) -> Dict[str, Any]:
        """
        Generate PPTX presentation from draft content.

        Args:
            input_data: DraftContent object with markdown content
            **kwargs:
                - brand_template: BrandTemplate for styling
                - slide_numbers: Add slide numbers (default: True)

        Returns:
            Dictionary with:
                - file_path: Path to generated PPTX file
                - success: Boolean indicating success
                - metadata: Additional information
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            self.logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return self._generate_mock_pptx(input_data, **kwargs)

        brand_template = kwargs.get("brand_template")
        slide_numbers = kwargs.get("slide_numbers", True)

        # Create presentation
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio)
        if brand_template and hasattr(brand_template, 'presentation_layout'):
            layout = brand_template.presentation_layout
            prs.slide_width = Inches(layout.slide_width)
            prs.slide_height = Inches(layout.slide_height)
        else:
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

        # Add title slide
        self._add_title_slide(prs, input_data, brand_template)

        # Parse content and create slides
        self._parse_content_to_slides(prs, input_data.content, brand_template)

        # Add closing slide
        self._add_closing_slide(prs, brand_template)

        # Generate filename and save
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{input_data.content_type.value}_{timestamp}.pptx"
        file_path = os.path.join(self.output_dir, filename)

        prs.save(file_path)
        self.logger.info(f"Generated PPTX: {file_path}")

        return {
            "file_path": file_path,
            "success": True,
            "metadata": {
                "word_count": input_data.word_count,
                "slides": len(prs.slides),
                "brand_template": brand_template.name if brand_template else None,
                "created_at": datetime.now().isoformat()
            }
        }

    def _add_title_slide(self, prs: Any, draft: DraftContent, brand_template: Any = None):
        """Add title slide to presentation."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        # Use blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = draft.content_type.value.replace('_', ' ').title()

        # Style title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True

        if brand_template:
            r, g, b = self._hex_to_rgb(brand_template.colors.primary)
            p.font.color.rgb = RGBColor(r, g, b)

            # Add company name subtitle
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(0.75)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = brand_template.company_name

            p = subtitle_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(24)
            r, g, b = self._hex_to_rgb(brand_template.colors.secondary)
            p.font.color.rgb = RGBColor(r, g, b)

        # Apply background color if brand template
        if brand_template:
            self._apply_slide_background(slide, brand_template)

    def _add_closing_slide(self, prs: Any, brand_template: Any = None):
        """Add closing/thank you slide."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        if not brand_template:
            return

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Thank you message
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Thank You"

        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True

        r, g, b = self._hex_to_rgb(brand_template.colors.primary)
        p.font.color.rgb = RGBColor(r, g, b)

        # Contact info
        contact_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(8), Inches(1.5)
        )
        contact_frame = contact_box.text_frame

        lines = [brand_template.company_name]
        if brand_template.website:
            lines.append(brand_template.website)

        contact_frame.text = '\n'.join(lines)

        for p in contact_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)
            r, g, b = self._hex_to_rgb(brand_template.colors.text)
            p.font.color.rgb = RGBColor(r, g, b)

        self._apply_slide_background(slide, brand_template)

    def _parse_content_to_slides(self, prs: Any, content: str, brand_template: Any = None):
        """Convert markdown content to slides."""
        lines = content.split('\n')
        current_slide = None
        current_content = []
        slide_title = None

        for line in lines:
            stripped = line.strip()

            # H1 = Section divider slide
            if stripped.startswith('# '):
                # Finalize previous slide if exists
                if current_slide and (slide_title or current_content):
                    self._finalize_content_slide(current_slide, slide_title, current_content, brand_template)

                # Create section slide
                section_title = stripped[2:].strip()
                self._add_section_slide(prs, section_title, brand_template)

                current_slide = None
                slide_title = None
                current_content = []

            # H2 = New content slide
            elif stripped.startswith('## '):
                # Finalize previous slide if exists
                if current_slide and (slide_title or current_content):
                    self._finalize_content_slide(current_slide, slide_title, current_content, brand_template)

                # Start new content slide
                slide_title = stripped[3:].strip()
                current_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                current_content = []

            # H3 or H4 = Sub-bullet or content
            elif stripped.startswith('### ') or stripped.startswith('#### '):
                text = stripped.lstrip('#').strip()
                if current_content:  # Add as sub-bullet
                    current_content.append(('sub', text))
                else:  # Add as regular bullet
                    current_content.append(('bullet', text))

            # Bullet points
            elif stripped.startswith('- ') or stripped.startswith('* '):
                text = stripped[2:].strip()
                current_content.append(('bullet', text))

            # Regular paragraph (if slide exists)
            elif stripped and current_slide:
                current_content.append(('bullet', stripped))

        # Finalize last slide
        if current_slide and (slide_title or current_content):
            self._finalize_content_slide(current_slide, slide_title, current_content, brand_template)

    def _add_section_slide(self, prs: Any, title: str, brand_template: Any = None):
        """Add a section divider slide."""
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add section title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title

        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True

        if brand_template:
            r, g, b = self._hex_to_rgb(brand_template.colors.primary)
            p.font.color.rgb = RGBColor(r, g, b)
            self._apply_slide_background(slide, brand_template)

    def _finalize_content_slide(self, slide: Any, title: Optional[str],
                                content: List, brand_template: Any = None):
        """Add title and content to a slide."""
        from pptx.util import Pt, Inches
        from pptx.enum.text import PP_ALIGN

        # Add title if provided
        y_position = Inches(0.5)
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), y_position, Inches(9), Inches(0.75)
            )
            title_frame = title_box.text_frame
            title_frame.text = title

            p = title_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

            if brand_template:
                r, g, b = self._hex_to_rgb(brand_template.colors.primary)
                p.font.color.rgb = RGBColor(r, g, b)

            y_position = Inches(1.5)

        # Add content
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), y_position, Inches(8.5), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, (item_type, text) in enumerate(content):
                if i > 0:
                    content_frame.add_paragraph()

                p = content_frame.paragraphs[-1]
                p.text = text
                p.level = 1 if item_type == 'sub' else 0
                p.font.size = Pt(16 if item_type == 'sub' else 18)

                if brand_template:
                    r, g, b = self._hex_to_rgb(brand_template.colors.text)
                    p.font.color.rgb = RGBColor(r, g, b)

        # Apply background
        if brand_template:
            self._apply_slide_background(slide, brand_template)

    def _apply_slide_background(self, slide: Any, brand_template: Any):
        """Apply brand background color to slide."""
        try:
            background = slide.background
            fill = background.fill
            fill.solid()

            r, g, b = self._hex_to_rgb(brand_template.colors.background)
            fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            self.logger.warning(f"Could not apply slide background: {e}")

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def _generate_mock_pptx(self, draft: DraftContent, **kwargs) -> Dict[str, Any]:
        """
        Generate a mock PPTX file when python-pptx is not available.
        Creates a simple text file with .pptx extension for development.
        """
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{draft.content_type.value}_{timestamp}.pptx"
        file_path = os.path.join(self.output_dir, filename)

        brand_template = kwargs.get("brand_template")

        # Create a simple text representation
        content = f"""PPTX Presentation (Mock)
{'=' * 50}

Content Type: {draft.content_type.value}
Word Count: {draft.word_count}
Created: {datetime.now().isoformat()}
"""

        if brand_template:
            content += f"""Brand Template: {brand_template.name}
Company: {brand_template.company_name}
"""

        content += f"""
{'=' * 50}

SLIDES:

Slide 1: Title Slide
- {draft.content_type.value.replace('_', ' ').title()}

{draft.content}

Closing Slide: Thank You
"""

        if brand_template:
            content += f"- {brand_template.company_name}\n"
            if brand_template.website:
                content += f"- {brand_template.website}\n"

        content += f"""
{'=' * 50}
Note: This is a mock PPTX file. Install python-pptx for real PPTX generation.
pip install python-pptx
"""

        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)

        self.logger.info(f"Generated mock PPTX: {file_path}")

        return {
            "file_path": file_path,
            "success": True,
            "metadata": {
                "word_count": draft.word_count,
                "mock": True,
                "created_at": datetime.now().isoformat()
            }
        }

    def validate_requirements(self) -> tuple[bool, list[str]]:
        """Validate that required dependencies are available."""
        missing = []

        try:
            import pptx
        except ImportError:
            missing.append("python-pptx")

        return len(missing) == 0, missing
